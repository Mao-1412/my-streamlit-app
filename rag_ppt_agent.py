# ======================
# LangChain + FAISS対応 完全版（修正版）
# ======================
import os
import re
import textwrap
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import streamlit as st
import openai

# --- LangChain追加 ---
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.docstore.document import Document
from langchain.schema import AIMessage, HumanMessage
from langchain.chat_models import ChatOpenAI

# Streamlit キャッシュをクリア
st.cache_data.clear()
st.cache_resource.clear()

# -----------------------
# BASE_PATH設定（Cloud/ローカル共通）
# -----------------------
if "STREAMLIT_SERVER" in os.environ:
    BASE_PATH = "."  # Streamlit Cloud 上
else:
    BASE_PATH = r"C:\rag_poc_2"  # ローカル環境

DATA_PATH = os.path.join(BASE_PATH, "data")
OUTPUT_PATH = os.path.join(BASE_PATH, "output")
VECTOR_PATH = os.path.join(BASE_PATH, "vectorstore")

# ディレクトリ作成（data以外）
os.makedirs(OUTPUT_PATH, exist_ok=True)
os.makedirs(VECTOR_PATH, exist_ok=True)

# dataフォルダの存在確認
if not os.path.exists(DATA_PATH) or not os.path.isdir(DATA_PATH):
    st.error(f"データフォルダが存在しません: {DATA_PATH}")
    data_files = []
else:
    data_files = os.listdir(DATA_PATH)
    print("Files in data folder:", data_files)

# OpenAI APIキー
openai.api_key = st.secrets.get("OPENAI_API_KEY")

# 日本語フォント設定
from matplotlib import rcParams
rcParams['font.family'] = 'MS Gothic'
sns.set(font='MS Gothic')


# -----------------------
# データ読み込み
# -----------------------
@st.cache_data
def load_data():
    try:
        pos_df = pd.read_excel(os.path.join(DATA_PATH, "POSデータ.xlsx"))
        delivery_df = pd.read_excel(os.path.join(DATA_PATH, "納品実績報告.xlsx"))
        store_df = pd.read_excel(os.path.join(DATA_PATH, "店舗情報.xlsx"))
        merch_df = pd.read_excel(os.path.join(DATA_PATH, "市場_POS実績比較.xlsx"))
        market_df = pd.read_excel(os.path.join(DATA_PATH, "市場動向_商品政策.xlsx"))
        store_display_df = pd.read_excel(os.path.join(DATA_PATH, "店だしデータ.xlsx"))
        client_df = pd.read_excel(os.path.join(DATA_PATH, "得意先情報.xlsx"))
        return pos_df, delivery_df, store_df, merch_df, market_df, store_display_df, client_df
    except FileNotFoundError as e:
        st.error(f"Excelファイルが見つかりません: {e}")
        # 空のDataFrameを返して以降の処理で落ちないようにする
        empty_df = pd.DataFrame()
        return (empty_df,) * 7

pos_df, delivery_df, store_df, merch_df, market_df, store_display_df, client_df = load_data()


# -----------------------
# RAG構築部分
# -----------------------
@st.cache_resource
def build_vectorstore():
    all_texts = []
    for df_name, df in {
        "POS": pos_df,
        "納品": delivery_df,
        "店舗": store_df,
        "市場比較": merch_df,
        "市場動向": market_df,
        "店だし": store_display_df,
        "得意先": client_df
    }.items():
        text_data = df.to_csv(index=False)
        doc = Document(page_content=text_data, metadata={"source": df_name})
        all_texts.append(doc)

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=150)
    docs = text_splitter.split_documents(all_texts)

    embeddings = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=openai.api_key)
    vectordb = FAISS.from_documents(docs, embedding=embeddings, persist_directory=VECTOR_PATH)
    vectordb.persist()
    return vectordb

vectorstore = build_vectorstore()

# -----------------------
# GPT修正関数（RAG対応）
# -----------------------
def refine_text_with_gpt(original_text, instruction):
    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    related_docs = retriever.get_relevant_documents(instruction)
    context_text = "\n".join([d.page_content[:500] for d in related_docs])

    prompt = f"""
    以下の関連データを参考に、次の文章を指示に従って修正してください。

    【関連データ】
    {context_text}

    【修正対象文章】
    {original_text}

    【修正指示】
    {instruction}
    """

    try:
        llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
        response = llm.invoke([HumanMessage(content=prompt)])
        return response.content.strip()
    except Exception as e:
        st.error(f"RAG修正に失敗しました: {e}")
        return original_text

# -----------------------
# 総括生成関数（RAG対応）
# -----------------------
def generate_summary_block(latest_blocks):
    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    related_docs = retriever.get_relevant_documents("営業提案の総括生成")
    context_text = "\n".join([d.page_content[:500] for d in related_docs])

    text_to_summarize = "\n\n".join([
        latest_blocks.get("【販売数量分析】", ""),
        latest_blocks.get("【商品提案】", ""),
        latest_blocks.get("【在庫管理】", "")
    ])

    if not any([latest_blocks.get(k) for k in ["【販売数量分析】","【商品提案】","【在庫管理】"]]):
        return "総括内容がありません。"

    prompt = f"""
    以下の関連データを踏まえて、次の内容を200字程度で営業提案用の総括文にまとめてください。

    【関連データ】
    {context_text}

    【要約対象】
    {text_to_summarize}
    """

    try:
        llm = ChatOpenAI(model="gpt-4o-mini", temperature=0)
        response = llm.invoke([HumanMessage(content=prompt)])
        return response.content.strip()
    except Exception as e:
        st.error(f"総括生成に失敗しました: {e}")
        return "総括の自動生成に失敗しました。"


# -----------------------
# PPT用フォント設定関数
# -----------------------
def set_font_for_text_frame(tf, font_name="Meiryo UI", font_size_pt=14, font_color=(0,0,0)):
    """
    pptx の TextFrame または Paragraph に対してフォント設定を統一する
    """
    if hasattr(tf, "paragraphs"):
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size_pt)
                run.font.color.rgb = RGBColor(*font_color)
    else:
        for run in tf.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = RGBColor(*font_color)

# -----------------------
# GPT提案文 自動生成
# -----------------------
def generate_gpt_proposal(client_name):
    prompt = f"""
    あなたはプロの営業コンサルタントです。
    得意先「{client_name}」に向けて、提案資料に掲載する自然文の提案内容を作成してください。
    トピック構成は以下の4つを含めてください。
    1. 新商品・新施策の提案
    2. 販促・デジタルマーケティング戦略
    3. オペレーション改善案
    4. 今後の成長方向性
    """
    try:
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}],
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"提案文生成に失敗しました: {e}")
        return "提案文の自動生成に失敗しました。"

# -----------------------
# セッション初期化
# -----------------------
if 'proposal_blocks' not in st.session_state:
    st.session_state['proposal_blocks'] = {}

if 'auto_generated_text' not in st.session_state:
    st.session_state['auto_generated_text'] = {}

def fix_proposal_block(block_key, auto_generated_text, instruction):
    """
    GPTにブロック文章を修正させる。
    文頭にブロック名は付けず、merge時に付与する
    """
    current_text = st.session_state['proposal_blocks'].get(block_key, auto_generated_text.get(block_key, ""))

    refined_text = refine_text_with_gpt(current_text, instruction)
    
    if refined_text:
        refined_text = refined_text.strip()
        # ここで「修正された文章:」を削除
        refined_text = re.sub(r'^(修正後の文章:|修正された文章:)\s*', '', refined_text)
        st.session_state['proposal_blocks'][block_key] = refined_text

    return refined_text


def merge_blocks_with_session(auto_generated_text):
    block_order = ["【販売数量分析】", "【商品提案】", "【在庫管理】", "【総括】"]
    final_lines = []
    for key in block_order:
        text = st.session_state.get('proposal_blocks', {}).get(key) or auto_generated_text.get(key)
        if text:
            # 「修正後の文章:」を削除
            text = re.sub(r'^修正後の文章:\s*', '', text)
            # ブロック名が含まれていない場合のみ追加
            if not text.startswith(key):
                final_lines.append(f"{key}\n{text.strip()}")
            else:
                final_lines.append(text.strip())
    return "\n\n".join(final_lines)

# -----------------------
# 売上チャート作成
# -----------------------
def create_sales_chart(series, chart_type="bar", title="売上チャート"):
    """
    series: pd.Series, index=日付, values=数量
    chart_type: "bar" または "line"
    return: BytesIO に保存した画像
    """
    fig, ax = plt.subplots(figsize=(6,4))
    if chart_type == "bar":
        series.plot(kind="bar", ax=ax, color="skyblue")
    else:
        series.plot(kind="line", ax=ax, marker='o', color="skyblue")
    ax.set_title(title)
    ax.set_ylabel("販売数量")
    plt.xticks(rotation=45)
    plt.tight_layout()

    buf = BytesIO()
    plt.savefig(buf, format='png')
    plt.close(fig)
    buf.seek(0)
    return buf

# -----------------------
# PPT生成
# -----------------------
def generate_ppt(client_name, client_id, proposal_text=None):
    prs = Presentation()

        # --- 修正保持対応 ---
    # 自動生成ではなく、修正済みテキストを優先して読み込む
    current_blocks = st.session_state.get('proposal_blocks', {}).copy()
    auto_generated_text = st.session_state.get('auto_generated_text', {}).copy()

    # 「修正済み」を常に優先してマージ
    for k, v in current_blocks.items():
        if v and v.strip():
            auto_generated_text[k] = v

    
    def preprocess_proposal_text(text):
        text = re.sub(r'[\r\v]+', '', text)
        # 「】」の後に改行を追加。ただし既に改行がある場合は重複しない
        text = re.sub(r'】(?!\n)', '】\n', text)
        lines = []
        for part in text.split("\n"):
            part = part.strip()
            if part:
                wrapped = textwrap.wrap(part, width=42)
                lines.extend(wrapped)
        return lines
    
    def add_slide(title, contents=None, img_bytes=None, is_proposal=False):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        try:
            slide.shapes.title.text = title
            set_font_for_text_frame(slide.shapes.title.text_frame, font_name="Meiryo UI", font_size_pt=18)
        except Exception:
            pass

        if contents:
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.clear()

            for content in contents:
                paragraph = tf.add_paragraph()
                paragraph.text = content
                set_font_for_text_frame(
                    tf, 
                    font_name="Meiryo UI", 
                    font_size_pt=14 if is_proposal else 16
                )

        if img_bytes:
            left, top = Inches(0.5), Inches(1.5)
            slide.shapes.add_picture(img_bytes, left, top, width=Inches(9), height=Inches(5))

        return slide

    # -----------------------
    # アジェンダ
    # -----------------------
    agenda_lines = [
        "1. 現状分析（売上・納品・市場動向）",
        "2. 店頭状況と販促状況",
        "3. 市場分析・競合比較",
        "4. 得意先分析・条件整理",
        "5. 提案内容",
        "6. 効果予測と次のアクション"
    ]
    add_slide(f"{client_name} 向け提案資料", agenda_lines)

    # -----------------------
    # 2. 現状分析（売上・納品状況）
    # -----------------------
    client_stores = store_df[store_df['得意先ID'] == client_id]['店舗ID'].tolist()
    pos_client = pos_df[pos_df['店舗ID'].isin(client_stores)]

    # 追加：市場POS比較データをエリア別に取得
    market_pos = merch_df[merch_df['エリア'].isin(store_df[store_df['得意先ID']==client_id]['エリア'])]

    img_buf = None
    if not pos_client.empty:
        pos_client_group = pos_client.groupby('販売日')['販売数量'].sum()
        img_buf = create_sales_chart(pos_client_group, chart_type="bar", title="月別販売数量")

    delivery_client = delivery_df[delivery_df['得意先ID']==client_id]
    if not delivery_client.empty:
        delivery_client['欠品率'] = delivery_client['欠品数量'] / delivery_client['受注数量'] * 100
        delivery_client['返品率'] = delivery_client['返品数量'] / delivery_client['受注数量'] * 100
        delivery_lines = [
            f"欠品率: {delivery_client['欠品率'].mean():.1f}%",
            f"返品率: {delivery_client['返品率'].mean():.1f}%"
        ]
    else:
        delivery_lines = ["納品データなし"]
    add_slide("現状分析（売上・納品状況）", delivery_lines, img_bytes=img_buf)

    # -----------------------
    # 3. 店頭状況と販促状況
    # -----------------------
    display_client = store_display_df[store_display_df['得意先ID']==client_id]
    fig, ax = plt.subplots(figsize=(6,4))
    if not display_client.empty:
        pivot_display = display_client.pivot(index='商品名', columns='店舗ID', values='店頭在庫数量')
        sns.heatmap(pivot_display.fillna(0), annot=True, fmt=".0f", cmap="YlGnBu", ax=ax)
        ax.set_title("店頭在庫ヒートマップ")
    else:
        ax.text(0.5, 0.5, "店だしデータなし", ha='center', va='center', fontsize=16)
        ax.set_title("店頭在庫ヒートマップ")
    plt.tight_layout()
    img_buf = BytesIO()
    plt.savefig(img_buf, format='png')
    plt.close(fig)
    img_buf.seek(0)
    add_slide("店頭状況と販促状況", ["店頭在庫状況ヒートマップ"], img_bytes=img_buf)

    # -----------------------
    # 4. 市場分析・競合比較
    # -----------------------
    fig, ax = plt.subplots(figsize=(6,4))
    if not market_pos.empty:
        ax.bar(market_pos['商品名'], market_pos['自社販売数量'], label="自社")
        ax.bar(market_pos['商品名'], market_pos['市場販売数量'], alpha=0.3, label="市場")
        ax.set_title("市場販売数量比較")
        plt.xticks(rotation=45)
        plt.legend()
    else:
        ax.text(0.5, 0.5, "市場データなし", ha='center', va='center', fontsize=16)
        ax.set_title("市場販売数量比較")
    plt.tight_layout()
    img_buf = BytesIO()
    plt.savefig(img_buf, format='png')
    plt.close(fig)
    img_buf.seek(0)
    add_slide("市場分析・競合比較", ["市場シェア比較"], img_bytes=img_buf)

    # -----------------------
    # 5. 得意先分析・条件整理
    # -----------------------
    client_info = client_df[client_df['得意先ID']==client_id]
    if not client_info.empty:
        client_info = client_info.iloc[0]
        lines = [
            f"店舗タイプ: {client_info['店舗タイプ']}",
            f"発注頻度: {client_info['納品頻度']}",
            f"特記事項: {client_info['特記事項']}"
        ]
    else:
        lines = ["得意先情報なし"]
    add_slide("得意先分析・条件整理", lines)
    # -----------------------
    # 6. データドリブン提案内容
    # -----------------------
    auto_generated_text = st.session_state.get('auto_generated_text', {})

    # 総括を必ずセッションにセット
    if "【総括】" not in st.session_state['proposal_blocks'] or not st.session_state['proposal_blocks']["【総括】"]:
        st.session_state['proposal_blocks']["【総括】"] = auto_generated_text.get("【総括】", "総括内容がありません。")


    store_info = store_df[store_df['得意先ID'] == client_id]

    # 【販売数量分析】
    if not pos_client.empty and not market_pos.empty:
        lines = []
        for _, row in market_pos.iterrows():
            diff = row.get("市場シェア（数量）", 0)
            competitiveness = row.get("価格競争力指数", 1.0)
            product = row.get("商品名", "商品不明")
            area = row.get("エリア", "地域不明")

            if diff < 1.0:
                lines.append(
                f"{area}エリアにおける「{product}」の市場シェアは{diff:.2%}で低水準です。"
                f"平均価格は競合と比べ{(1 - competitiveness)*100:.1f}%低く、価格競争力があります。"
                f"消費者により良く目に留まる陳列位置の改善や、週末販売の強化によって売上向上が期待できます。"
            )
            elif diff >= 1.0 and competitiveness < 0.99:
                lines.append(
                    f"【価格戦略分析】{product}はシェア好調（{diff:.2%}）ですが、単価が市場より低いため、"
                    f"利益率向上のために高価格帯ギフト仕様商品の追加を検討すると良いでしょう。"
                )
        auto_generated_text["【販売数量分析】"] = "\n".join(lines)
    else:
        auto_generated_text["【販売数量分析】"] = "販売数量分析に必要なデータが不足しています。"

    # 【商品提案】
    lines = []
    if not store_info.empty:
        for _, row in store_info.iterrows():
            store_type = row.get("店舗タイプ", "")
            store_name = row.get("店舗名", "")
            if store_type in ["スーパー", "ショッピングモール"]:
                lines.append(f"【販促提案】{store_name}では試食やイベント連動の販促が有効です。")
            elif store_type in ["専門店", "高級専門店"]:
                lines.append(f"【商品提案】{store_name}では高単価・ギフト対応商品の強化を推奨します。")
    else:
        lines.append("【商品提案】店舗情報が不足しているため、一般的な提案として季節限定商品や地域特化型商品の導入を検討してください。")
    auto_generated_text["【商品提案】"] = "\n".join(lines)

    # 【在庫管理】
    if not delivery_client.empty:
        if delivery_client["欠品数量"].sum() > 0:
            auto_generated_text["【在庫管理】"] = "欠品が一部発生しているため、納品頻度の増加や在庫補充体制の強化を検討してください。"
        else:
            auto_generated_text["【在庫管理】"] = "欠品は発生しておらず、現行頻度で安定供給可能です。"
    else:
        auto_generated_text["【在庫管理】"] = "納品データが不足しているため、在庫状況の把握を強化する必要があります。"

    # -----------------------
    # 【総括】を動的生成
    # -----------------------
    # 最新のブロック情報を統合して必ず定義
    latest_blocks = {**st.session_state.get('auto_generated_text', {}), 
                    **st.session_state.get('proposal_blocks', {})}

    # いずれかのブロックに内容がある場合のみ総括を生成
    if any(latest_blocks.get(k) for k in ["【販売数量分析】", "【商品提案】", "【在庫管理】"]):
        auto_generated_text["【総括】"] = generate_summary_block(latest_blocks)
    else:
        auto_generated_text["【総括】"] = "総括内容がありません。"

    # セッションに反映
    st.session_state['auto_generated_text'] = auto_generated_text


    # ブロック統合とスライド生成（統合版）
    proposal_blocks = {**st.session_state.get('auto_generated_text', {}), **st.session_state.get('proposal_blocks', {})}
    proposal_blocks["【総括】"] = generate_summary_block(proposal_blocks)
    st.session_state['proposal_blocks'] = proposal_blocks

    # PPT用テキスト生成
    if proposal_text:  # 引数がある場合はそれを使用
        final_text = preprocess_proposal_text(proposal_text)
    else:  # 引数がない場合はセッションの最新ブロックをマージして使用
        merged_text = merge_blocks_with_session(st.session_state['proposal_blocks'])
        final_text = preprocess_proposal_text(merged_text)

    # セッションを更新
    st.session_state['auto_generated_text'] = auto_generated_text

    # ローカル変数を直接使用してマージ（ここがポイント！）
    # 統合ブロックをPPT用テキストに変換
    if proposal_text:  # 引数がある場合はそれを使用
        final_text = preprocess_proposal_text(proposal_text)
    else:  # 引数がない場合はセッションの最新ブロックをマージして使用
        merged_text = merge_blocks_with_session(st.session_state['proposal_blocks'])
        final_text = preprocess_proposal_text(merged_text)

    # ここで一度だけスライドを追加
    add_slide("提案内容", final_text, is_proposal=True)

    # 保存
    output_file = os.path.join(OUTPUT_PATH, f"{client_name}_proposal.pptx")
    prs.save(output_file)
    return output_file

# -----------------------
# Streamlit UI
# -----------------------
st.title("自然文PPT＋キャッチボール型提案資料生成（完全最終版）")

client_name = st.text_input("得意先名を入力", "")

if st.button("初回PPT生成"):
    client_row = client_df[client_df['得意先名']==client_name]
    if client_row.empty:
        st.error(f"{client_name} はデータに存在しません。")
    else:
        client_id = client_row.iloc[0]['得意先ID']
        ppt_file = generate_ppt(client_name, client_id)
        st.session_state['ppt_file'] = ppt_file
        st.success(f"PPTを生成しました: {ppt_file}")
        with open(ppt_file, "rb") as f:
            st.download_button("PPTをダウンロード", f, file_name=os.path.basename(ppt_file))


# -----------------------
# ブロック単位修正UI
# -----------------------
block_to_fix = st.selectbox(
    "修正対象ブロックを選択",
    ["【販売数量分析】","【商品提案】","【在庫管理】","【総括】"]
)
instruction = st.text_area(f"{block_to_fix} の修正指示を入力")

if st.button("ブロック修正＆再生成"):
    if not instruction:
        st.error("修正指示を入力してください。")
    else:
        client_row = client_df[client_df['得意先名']==client_name]
        if client_row.empty:
            st.error(f"{client_name} はデータに存在しません。")
        else:
            client_id = client_row.iloc[0]['得意先ID']

            # ① 前回修正済みのブロック文章をセッションから取得
            auto_generated_text = st.session_state.get('auto_generated_text', {})
            proposal_blocks = st.session_state.get('proposal_blocks', {})

            # ② GPTでブロック修正
            refined_text = fix_proposal_block(
                block_key=block_to_fix,
                auto_generated_text={**auto_generated_text, **proposal_blocks},  # 最新状態を渡す
                instruction=instruction
            )

            # ③ 修正結果をセッションに反映（前回の修正は維持）
            st.session_state['proposal_blocks'][block_to_fix] = refined_text

            # ④ 総括も自動更新
            latest_blocks = {**st.session_state.get('auto_generated_text', {}), **st.session_state.get('proposal_blocks', {})}
            st.session_state['proposal_blocks']["【総括】"] = generate_summary_block(latest_blocks)

            # ⑤ PPT生成時に proposal_blocks と auto_generated_text を統合して使用
            ppt_file = generate_ppt(client_name, client_id, proposal_text=None)
            st.session_state['ppt_file'] = ppt_file

            st.success("修正版PPTを生成しました。")
            
            # 不要な「修正後の文章:」を削除して表示
            clean_text = re.sub(r'^修正後の文章:\s*', '', refined_text)
            st.text_area("修正版ブロック文章", clean_text, height=300)
            
            with open(ppt_file, "rb") as f:
                st.download_button(
                    "修正版PPTをダウンロード",
                    f,
                    file_name=os.path.basename(ppt_file)
                )













